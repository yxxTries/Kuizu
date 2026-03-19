"""Extract slide content from .pptx and .pdf files."""
from pathlib import Path
from app.models import SlideChunk

def parse_pptx(path: Path) -> list[SlideChunk]:
    from pptx import Presentation
    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:
                title = shape.text.strip()
            else:
                bullets.extend(
                    p.text.strip()
                    for p in shape.text_frame.paragraphs
                    if p.text.strip()
                )
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if title or bullets:
            chunks.append(SlideChunk(
                slide_number=i,
                title=title,
                content="\n".join(bullets),
                notes=notes or None,
            ))
    return chunks

def parse_pdf(path: Path) -> list[SlideChunk]:
    import pdfplumber
    chunks = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            lines = [l.strip() for l in text.splitlines() if l.strip()]
            title = lines[0] if lines else f"Slide {i}"
            content = "\n".join(lines[1:])
            if content:
                chunks.append(SlideChunk(
                    slide_number=i,
                    title=title,
                    content=content,
                ))
    return chunks
